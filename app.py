import streamlit as st
import fitz  # PyMuPDF for PDF handling
import cohere
from PIL import Image
import io
from docx import Document
from pptx import Presentation
import textwrap
from gtts import gTTS
import os
import tempfile

# Initialize Cohere client
COHERE_API_KEY = 'jwI4BsEEPcXezrMzjZlJyqwdgqPK0RZBpVU95HmP'
co = cohere.Client(COHERE_API_KEY)

# Constants
TOKEN_LIMIT = 4081  # Adjust based on the API's token limit
CHUNK_SIZE = 2000  # Size of each text chunk to avoid exceeding token limits

# Extract text and images from PDF
def extract_text_and_images_from_pdf(file):
    text = ""
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    
    # Extract text using PyMuPDF
    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        text += page.get_text()

    # Extract images using PyMuPDF
    images_text = ""
    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        for img_index, img in enumerate(page.get_images(full=True)):
            base_image = fitz.Pixmap(pdf_document, img[0])
            if base_image.n < 5:  # GRAYSCALE or RGB image
                try:
                    img_data = base_image.tobytes("png")
                    pil_image = Image.open(io.BytesIO(img_data))
                    # images_text += pytesseract.image_to_string(pil_image)
                except ValueError:
                    # Unsupported color space, skip the image
                    continue
            else:
                # Handle CMYK or other color spaces if needed
                # For now, skip these images
                continue

    return text, images_text

# Extract text from Word document
def extract_text_from_docx(file):
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Extract text from PowerPoint presentation
def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Summarize a chunk of text using Cohere API
def summarize_chunk(chunk):
    response = co.generate(
        model='command',  # Replace with the correct model ID
        prompt=f"Summarize the following text:\n\n{chunk}\n\nSummary:",
        max_tokens=300  # Adjust as needed
    )
    return response.generations[0].text.strip()

# Chunk and summarize text
def summarize_text(text):
    chunks = textwrap.wrap(text, width=CHUNK_SIZE)  # Chunk text
    summaries = [summarize_chunk(chunk) for chunk in chunks]
    return " ".join(summaries)

# Answer questions using Cohere API with chunking
def answer_question(question, context):
    chunks = textwrap.wrap(context, width=CHUNK_SIZE)  # Chunk text
    answers = []
    
    for chunk in chunks:
        response = co.generate(
            model='command',  # Replace with the correct model ID
            prompt=f"Context: {chunk}\n\nQuestion: {question}\nAnswer:",
            max_tokens=300  # Adjust as needed
        )
        answers.append(response.generations[0].text.strip())
    
    return " ".join(answers)

# Generate additional context-related data
def generate_additional_context(context):
    chunks = textwrap.wrap(context, width=CHUNK_SIZE)  # Chunk text
    additional_contexts = []
    
    for chunk in chunks:
        response = co.generate(
            model='command',  # Replace with the correct model ID
            prompt=f"Context: {chunk}\n\nProvide additional context related to the document:",
            max_tokens=400  # Adjust as needed
        )
        additional_contexts.append(response.generations[0].text.strip())
    
    return " ".join(additional_contexts)

# Convert text to speech and save to a temporary file
def text_to_speech(text):
    tts = gTTS(text, lang='en')
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
    tts.save(temp_file.name)
    return temp_file.name

def main():
    st.markdown("""
        <style>
            .main {
                background: grey;
                padding: 20px;
                border-radius: 10px;
                box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);               
                backdrop-filter: blur(5px);
                
            }
            
            .title {
                
                font-family: Segoe ui, sans-serif; 
 
                background-clip: text; 
                -webkit-background-clip: text; 
                color: white;
                background-image: white; 
                letter-spacing: 0.5px; 
                font-size:30px;
                font-weight:900;
                height:60px;                            
            }
            
            
            .subheader {
                top: 60px;
                color: #242324;
                font-size: 24px;
                font-weight: bold;
            }
            .text {
                font-size: 16px;
                line-height: 1.5;
                color: #242324;
                
            }
            .text-container {
                border: 4px solid #261C30;
                border-radius: 10px;
                padding: 20px;
                background-color: #f8f8f8;
            }
            .button {
               background-color: #007bff;
               
               padding: 10px 20px;
               border: none;
               border-radius: 5px;
               cursor: pointer;
                
            }
            .button:hover {
                background-color: #007bff;
                color:
            }
            .st-emotion-cache-1qg05tj {
                font-size: 20px;
                color: rgb(0 0 0);
               
            }
            .st-emotion-cache-13k62yr {
                position: absolute;
                background: rgb(14, 17, 23);
                color: rgb(0, 0, 0);
                inset: 0px;
                color-scheme: dark;
                overflow: hidden;
            }
            .st-b6 {
               
               color: rgb(255 246 246);
              
               
            }
            .checkbox {
                cursor: pointer;
                color: rgb(255, 255, 255);
            }
            .st-emotion-cache-13k62yr {
                position: absolute;
                background: rgb(0, 0, 0);
                color: rgb(0, 0, 0);
                inset: 0px;
                color-scheme: dark;
                overflow: hidden;
            }
                .stTextInput > div {
                    width: 100% !important;
                    max-width: 100% !important;
                }
                .stTextInput input {
                    width: 100% !important;
                    max-width: 100% !important;
                }
                .stButton > button {
                    background-color: #000000; /* Green background */
                    color: white; /* White text */
                    border: none; /* Remove borders */
                    padding: 10px 20px; /* Add some padding */
                    text-align: center; /* Center the text */
                    text-decoration: none; /* Remove underline */
                    display: inline-block; /* Make the button inline */
                    font-size: 16px; /* Increase font size */
                    margin: 4px 2px; /* Add some margin */
                    cursor: pointer; /* Add a pointer cursor on hover */
                    border-radius: 4px; /* Rounded corners */
                }

                .st-cb label {
                    /* Your desired styles here */
                    color: #ffffff; /* Change the text color */
                    font-weight: bold; /* Make the text bold */
                }
                
                

           
            
            
            
        </style>
    """, unsafe_allow_html=True)

    # st.markdown('<div class="main">', unsafe_allow_html=True)
    st.markdown('<div class="title">Document Summarizer and Q&A</div>', unsafe_allow_html=True)

    # st.markdown("""
    #     <style>
    #         .stSelectbox{
    #             font-size: 600px;
    #             font-weight: bold;
    #             color: #4CAF50;
    #             margin-bottom: 10px;
    #         }
    #     </style>
    #     """, unsafe_allow_html=True)
    
    st.markdown('<div class="st-emotion-cache-1qg05tj "></div>', unsafe_allow_html=True)

    file_type = st.selectbox("Select file type", ["Select", "PDF", "Word", "PowerPoint"])

    if file_type == "Select":
        st.write("Please select a document type.")
        return

    if file_type == "PDF":
        file_types = ["pdf"]
    elif file_type == "Word":
        file_types = ["docx"]
    elif file_type == "PowerPoint":
        file_types = ["pptx"]
    else:
        st.write("Invalid file type selected.")
        return

    uploaded_file = st.file_uploader(f"Upload a {file_type} file", type=file_types)

    if uploaded_file is not None:
        if file_type == "PDF":
            text, images_text = extract_text_and_images_from_pdf(uploaded_file)
        elif file_type == "Word":
            text = extract_text_from_docx(uploaded_file)
            images_text = ""  # No image extraction for Word
        elif file_type == "PowerPoint":
            text = extract_text_from_pptx(uploaded_file)
            images_text = ""  # No image extraction for PowerPoint

        combined_text = text + " " + images_text

        if combined_text:
            # st.markdown('<div class="st-cb"></div>', unsafe_allow_html=True)
            
            
            
            show_text = st.checkbox("Show Document Text", value=False)
            

            if show_text:
                st.markdown('<div class="subheader">Extracted Text</div>', unsafe_allow_html=True)

                st.markdown(f'<div class ="text-container"><div class="text">{combined_text}</div></div>', unsafe_allow_html=True)
                print("\n")
            
            # st.markdown('<div class="st-emotion-cache-15hul6a"></div>', unsafe_allow_html=True)
            if st.button("Summarize Text", key="summarize"):
                
                # st.markdown('<div class="subheader">Summarized Text</div>', unsafe_allow_html=True)
                summary = summarize_text(combined_text)
                st.markdown(f'<div class ="text-container"><div class="text">{summary}</div></div>', unsafe_allow_html=True)

                # Generate and play speech for the summary
                if summary:
                    audio_file = text_to_speech(summary)
                    st.audio(audio_file)

            st.markdown('<div class="subheader">Ask a Question</div>', unsafe_allow_html=True)
            question = st.text_input("Enter your question:", key="question")
            
            
            if question:
                # Generate additional context related to the document
                additional_context = generate_additional_context(combined_text)
                
                # Combine additional context with the extracted text
                full_context = combined_text + " " + additional_context
                answer = answer_question(question, full_context)
                st.markdown(f'<div class ="text-container"><div class="text">{answer}</div></div>', unsafe_allow_html=True)

                # Generate and play speech for the answer
                if answer:
                    audio_file = text_to_speech(answer)
                    st.audio(audio_file)
        else:
            st.error("Could not extract text from the uploaded file.")
    st.markdown('</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
